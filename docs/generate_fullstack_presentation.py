#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Data Warehouse Full Stack Workshop.
Run: python generate_fullstack_presentation.py
Output: Data_Warehouse_Full_Stack.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Helper to create RGB color
def RgbColor(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# Color scheme - Professional dark theme
COLORS = {
    'bg_dark': RgbColor(15, 23, 42),       # Slate 900
    'bg_card': RgbColor(30, 41, 59),       # Slate 800
    'primary': RgbColor(59, 130, 246),     # Blue 500
    'accent': RgbColor(16, 185, 129),      # Emerald 500
    'text': RgbColor(248, 250, 252),       # Slate 50
    'text_muted': RgbColor(148, 163, 184), # Slate 400
    'danger': RgbColor(239, 68, 68),       # Red 500
    'warning': RgbColor(245, 158, 11),     # Amber 500
    'purple': RgbColor(139, 92, 246),      # Violet 500
    'cyan': RgbColor(6, 182, 212),         # Cyan 500
}


def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_agenda_slide(prs)
    
    # Part 1: Introduction
    add_section_slide(prs, "Part 1", "Introduction & Business Value", "5 min")
    add_modern_stack_slide(prs)
    add_stack_overview_slide(prs)
    add_tech_stack_slide(prs)
    add_business_value_slide(prs)
    
    # Part 2: Architecture
    add_section_slide(prs, "Part 2", "Architecture Overview", "10 min")
    add_high_level_arch_slide(prs)
    add_services_slide(prs)
    add_data_flow_slide(prs)
    
    # Part 3: Data Pipeline
    add_section_slide(prs, "Part 3", "Data Pipeline: Source → Transform", "10 min")
    add_adventureworks_slide(prs)
    add_airbyte_slide(prs)
    add_dbt_overview_slide(prs)
    add_dbt_layers_slide(prs)
    add_dbt_marts_slide(prs)
    
    # Part 4: Analytics Layer
    add_section_slide(prs, "Part 4", "Analytics: Dashboards & AI", "10 min")
    add_streamlit_slide(prs)
    add_ai_assistant_slide(prs)
    add_ai_architecture_slide(prs)
    add_ai_security_slide(prs)
    
    # Part 5: Security
    add_section_slide(prs, "Part 5", "Security & Access Control", "10 min")
    add_why_policy_slide(prs)
    add_opa_opal_slide(prs)
    add_roles_slide(prs)
    add_access_matrix_slide(prs)
    
    # Part 6: Operations
    add_section_slide(prs, "Part 6", "Live Demo & Operations", "10 min")
    add_start_stack_slide(prs)
    add_commands_slide(prs)
    
    # Part 7: Wrap-up
    add_section_slide(prs, "Part 7", "Q&A and Next Steps", "5 min")
    add_takeaways_slide(prs)
    add_best_practices_slide(prs)
    add_next_steps_slide(prs)
    add_resources_slide(prs)
    add_thank_you_slide(prs)
    
    return prs


def add_slide(prs, layout_index=6):
    """Add a blank slide."""
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)


def add_bg(slide, prs, color=None):
    """Add background to slide."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color or COLORS['bg_dark']
    bg.line.fill.background()
    return bg


def add_title_text(slide, text, left, top, width, height, font_size=44, bold=True, color=None):
    """Add a title text box."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text']
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_body_text(slide, text, left, top, width, height, font_size=18, color=None):
    """Add body text."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS['text_muted']
    return shape


def add_code_box(slide, code, left, top, width, height, font_size=11):
    """Add a code box with monospace font."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(30, 30, 46)
    shape.line.color.rgb = RgbColor(69, 71, 90)
    
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    
    for i, line in enumerate(code.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Consolas'
        p.font.color.rgb = RgbColor(166, 227, 161)
    
    return shape


def add_title_slide(prs):
    """Add the title slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    # Icons
    add_title_text(slide, "📊", 5.8, 1.2, 2, 1.2, font_size=80, bold=False)
    
    # Main title
    add_title_text(slide, "Modern Data Warehouse", 2, 2.5, 9, 1, font_size=52)
    add_title_text(slide, "Full Stack Architecture", 2, 3.4, 9, 0.8, font_size=40, 
                   color=COLORS['primary'])
    
    # Subtitle
    add_body_text(slide, "Complete Architecture & Implementation Workshop", 2, 4.6, 9, 0.5, 
                  font_size=22, color=COLORS['text_muted'])
    
    # Tech badges
    techs = ["PostgreSQL", "dbt", "Airbyte", "Streamlit", "AI/GPT-4", "OPA/OPAL"]
    for i, tech in enumerate(techs):
        x = 1.5 + i * 1.8
        add_body_text(slide, f"• {tech}", x, 5.8, 1.7, 0.4, font_size=14, color=COLORS['cyan'])


def add_agenda_slide(prs):
    """Add the agenda slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📋 Agenda (60 minutes)", 0.5, 0.4, 10, 0.8, font_size=36)
    
    agenda_items = [
        ("0:00 - 0:05", "Introduction & Business Value", "🎯"),
        ("0:05 - 0:15", "Architecture Overview", "🏗️"),
        ("0:15 - 0:25", "Data Pipeline: Source → Transform", "🔄"),
        ("0:25 - 0:35", "Analytics: Dashboards & AI", "📊"),
        ("0:35 - 0:45", "Security & Access Control", "🔐"),
        ("0:45 - 0:55", "Live Demo & Operations", "💻"),
        ("0:55 - 1:00", "Q&A and Next Steps", "❓"),
    ]
    
    for i, (time, topic, icon) in enumerate(agenda_items):
        y = 1.4 + i * 0.75
        add_body_text(slide, time, 1, y, 2, 0.5, font_size=16, color=COLORS['accent'])
        add_body_text(slide, f"{icon}  {topic}", 3.5, y, 9, 0.5, font_size=20, color=COLORS['text'])


def add_section_slide(prs, part, title, duration):
    """Add a section divider slide."""
    slide = add_slide(prs)
    add_bg(slide, prs, RgbColor(23, 37, 84))  # Darker blue
    
    add_body_text(slide, part, 0.5, 2.5, 12, 0.6, font_size=24, color=COLORS['primary'])
    add_title_text(slide, title, 0.5, 3.2, 12, 1, font_size=46)
    add_body_text(slide, f"⏱️ {duration}", 0.5, 4.5, 3, 0.5, font_size=18, color=COLORS['text_muted'])


def add_modern_stack_slide(prs):
    """Add modern vs traditional slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔄 Traditional vs Modern Data Stack", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Traditional column
    add_title_text(slide, "Traditional", 1, 1.4, 5, 0.5, font_size=24, color=COLORS['danger'])
    traditional = [
        "❌ Monolithic ETL tools",
        "❌ On-premise servers",
        "❌ SQL-only analytics",
        "❌ Static reports",
        "❌ Manual access control",
    ]
    for i, item in enumerate(traditional):
        add_body_text(slide, item, 1.2, 2 + i * 0.6, 5, 0.4, font_size=16)
    
    # Modern column
    add_title_text(slide, "Modern", 7, 1.4, 5, 0.5, font_size=24, color=COLORS['accent'])
    modern = [
        "✅ Modular, best-of-breed tools",
        "✅ Cloud-native & containerized",
        "✅ SQL + AI/ML integration",
        "✅ Interactive dashboards",
        "✅ Policy-as-code authorization",
    ]
    for i, item in enumerate(modern):
        add_body_text(slide, item, 7.2, 2 + i * 0.6, 5, 0.4, font_size=16, color=COLORS['accent'])


def add_stack_overview_slide(prs):
    """Add stack overview diagram."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🏗️ Our Data Warehouse Stack", 0.5, 0.4, 12, 0.8, font_size=34)
    
    diagram = """
    ┌───────────────────────────────────────────────────────────────┐
    │                   Modern Data Warehouse Stack                  │
    ├───────────────────────────────────────────────────────────────┤
    │                                                               │
    │  ┌──────────┐    ┌──────────┐    ┌──────────┐               │
    │  │SQL Server│───►│ Airbyte  │───►│PostgreSQL│               │
    │  │Adventure │    │   (EL)   │    │  (DWH)   │               │
    │  │  Works   │    │          │    │          │               │
    │  └──────────┘    └──────────┘    └────┬─────┘               │
    │                                       │                      │
    │                              ┌────────▼────────┐             │
    │                              │       dbt       │             │
    │                              │  (Transform)    │             │
    │                              └────────┬────────┘             │
    │                                       │                      │
    │            ┌──────────────────────────┼──────────────────┐   │
    │            │                          │                  │   │
    │    ┌───────▼───────┐          ┌──────▼──────┐           │   │
    │    │   Streamlit   │          │  dbt Docs   │           │   │
    │    │  + AI Assist  │          │  (Catalog)  │           │   │
    │    └───────┬───────┘          └─────────────┘           │   │
    │            │                                             │   │
    │    ┌───────▼───────┐                                    │   │
    │    │   OPA/OPAL    │                                    │   │
    │    │  (Security)   │                                    │   │
    │    └───────────────┘                                    │   │
    └───────────────────────────────────────────────────────────────┘"""
    
    add_code_box(slide, diagram, 0.8, 1.1, 11.7, 5.8, font_size=9)


def add_tech_stack_slide(prs):
    """Add technology stack table."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🛠️ Technology Stack", 0.5, 0.4, 12, 0.8, font_size=34)
    
    stack = [
        ("🗄️", "Source", "SQL Server + AdventureWorks", "Enterprise sample data"),
        ("🔄", "Extract/Load", "Airbyte", "Data ingestion (300+ connectors)"),
        ("💾", "Storage", "PostgreSQL 16", "Data warehouse"),
        ("⚙️", "Transform", "dbt Core", "SQL-based ELT"),
        ("📊", "Visualization", "Streamlit", "Interactive dashboards"),
        ("🤖", "AI", "OpenAI GPT-4 / Claude", "Natural language queries"),
        ("🔐", "Security", "OPA + OPAL", "Policy-based access control"),
        ("🐳", "Deployment", "Docker Compose", "Container orchestration"),
    ]
    
    for i, (icon, layer, tech, purpose) in enumerate(stack):
        y = 1.2 + i * 0.7
        add_title_text(slide, icon, 0.6, y, 0.5, 0.5, font_size=20, bold=False)
        add_body_text(slide, layer, 1.2, y, 1.8, 0.4, font_size=14, color=COLORS['primary'])
        add_title_text(slide, tech, 3.2, y, 4, 0.4, font_size=16, color=COLORS['text'])
        add_body_text(slide, purpose, 7.5, y, 5, 0.4, font_size=13)


def add_business_value_slide(prs):
    """Add business value slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "💡 Business Value", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Data Teams
    add_title_text(slide, "For Data Teams", 0.8, 1.3, 4, 0.5, font_size=20, color=COLORS['primary'])
    data_team = ["⚡ Fast Setup - hours, not weeks", "🔄 Modern ELT patterns", "📊 Self-Service BI"]
    for i, item in enumerate(data_team):
        add_body_text(slide, item, 1, 1.9 + i * 0.5, 4, 0.4, font_size=14)
    
    # Business Users
    add_title_text(slide, "For Business Users", 0.8, 3.7, 4, 0.5, font_size=20, color=COLORS['accent'])
    business = ["🤖 AI Analytics - ask in English", "📈 Real-time Insights", "🔐 Secure Access"]
    for i, item in enumerate(business):
        add_body_text(slide, item, 1, 4.3 + i * 0.5, 4, 0.4, font_size=14, color=COLORS['accent'])
    
    # Engineering
    add_title_text(slide, "For IT/Engineering", 6.5, 1.3, 5, 0.5, font_size=20, color=COLORS['purple'])
    eng = ["🐳 Containerized - easy deploy", "📝 Code-based - version controlled", "🔌 Modular - swap components"]
    for i, item in enumerate(eng):
        add_body_text(slide, item, 6.7, 1.9 + i * 0.5, 5, 0.4, font_size=14)


def add_high_level_arch_slide(prs):
    """Add high-level architecture slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🏛️ High-Level Architecture", 0.5, 0.4, 12, 0.8, font_size=34)
    
    arch = """
                         Data Sources
                    ┌────────────────────┐
                    │ SQL Server, APIs.. │
                    └─────────┬──────────┘
                              │
                              ▼
                    ┌────────────────────┐
                    │   Airbyte (EL)     │  ← Extract & Load
                    └─────────┬──────────┘
                              │
                              ▼
                    ┌────────────────────┐
                    │    PostgreSQL      │  ← Storage
                    └─────────┬──────────┘
                              │
                              ▼
                    ┌────────────────────┐
                    │   dbt (Transform)  │  ← staging → marts
                    └─────────┬──────────┘
                              │
          ┌───────────────────┼───────────────────┐
          │                   │                   │
          ▼                   ▼                   ▼
    ┌───────────┐       ┌───────────┐       ┌───────────┐
    │ Streamlit │       │ dbt Docs  │       │ OPA/OPAL  │
    │ Dashboard │       │ (Catalog) │       │ (Security)│
    └───────────┘       └───────────┘       └───────────┘"""
    
    add_code_box(slide, arch, 1.5, 1.1, 10.3, 5.8, font_size=10)


def add_services_slide(prs):
    """Add services overview slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔌 Service Architecture", 0.5, 0.4, 12, 0.8, font_size=34)
    
    services = [
        ("📊", "PostgreSQL", "5432", "data_warehouse_postgres", "Data storage"),
        ("📈", "Streamlit", "8501", "data_warehouse_streamlit", "Dashboard UI"),
        ("📚", "dbt Docs", "8080", "data_warehouse_dbt_docs", "Data catalog"),
        ("🗄️", "SQL Server", "1433", "data_warehouse_sqlserver", "Source database"),
        ("🔄", "Airbyte", "8000", "airbyte-abctl-*", "Data ingestion"),
        ("🔐", "OPA", "8181", "opal-opa-standalone", "Policy engine"),
        ("🔐", "OPAL Server", "7002", "opal-server", "Policy admin"),
    ]
    
    # Headers
    add_body_text(slide, "Service", 1.5, 1.3, 2, 0.4, font_size=14, color=COLORS['primary'])
    add_body_text(slide, "Port", 4, 1.3, 1, 0.4, font_size=14, color=COLORS['primary'])
    add_body_text(slide, "Container", 5.5, 1.3, 3.5, 0.4, font_size=14, color=COLORS['primary'])
    add_body_text(slide, "Purpose", 9.5, 1.3, 3, 0.4, font_size=14, color=COLORS['primary'])
    
    for i, (icon, service, port, container, purpose) in enumerate(services):
        y = 1.8 + i * 0.65
        add_title_text(slide, icon, 0.8, y, 0.5, 0.4, font_size=16, bold=False)
        add_body_text(slide, service, 1.5, y, 2.3, 0.4, font_size=14, color=COLORS['text'])
        add_body_text(slide, port, 4, y, 1, 0.4, font_size=14, color=COLORS['cyan'])
        add_body_text(slide, container, 5.5, y, 3.5, 0.4, font_size=12)
        add_body_text(slide, purpose, 9.5, y, 3, 0.4, font_size=12)


def add_data_flow_slide(prs):
    """Add data flow slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔄 Data Flow", 0.5, 0.4, 12, 0.8, font_size=34)
    
    steps = [
        ("1", "Source Data", "SQL Server with AdventureWorks: HR, Production, Purchasing, Sales"),
        ("2", "Extract & Load", "Airbyte: Full refresh or incremental sync, schema auto-detection"),
        ("3", "Transform", "dbt: Raw → Staging → Intermediate → Marts"),
        ("4", "Consume", "Streamlit dashboards + AI Assistant"),
    ]
    
    for i, (num, title, desc) in enumerate(steps):
        y = 1.3 + i * 1.4
        
        # Number circle
        add_title_text(slide, num, 0.9, y, 0.6, 0.6, font_size=24, color=COLORS['primary'])
        
        # Title and description
        add_title_text(slide, title, 1.8, y, 10, 0.5, font_size=20, color=COLORS['text'])
        add_body_text(slide, desc, 1.8, y + 0.5, 10, 0.5, font_size=14)
        
        # Arrow (except last)
        if i < len(steps) - 1:
            add_body_text(slide, "↓", 1, y + 1, 0.5, 0.3, font_size=18, color=COLORS['accent'])


def add_adventureworks_slide(prs):
    """Add AdventureWorks slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📦 Data Source: AdventureWorks", 0.5, 0.4, 12, 0.8, font_size=34)
    add_body_text(slide, "Microsoft's sample enterprise database - a fictional bicycle company", 
                  0.5, 1.1, 12, 0.4, font_size=18, color=COLORS['text'])
    
    schemas = [
        ("👥", "HumanResources", "Employee, Department, JobCandidate"),
        ("🏭", "Production", "Product, WorkOrder, BillOfMaterials"),
        ("📦", "Purchasing", "Vendor, PurchaseOrderHeader"),
        ("💰", "Sales", "Customer, SalesOrderHeader, Territory"),
        ("👤", "Person", "Person, Address, EmailAddress"),
    ]
    
    for i, (icon, schema, tables) in enumerate(schemas):
        y = 1.8 + i * 0.9
        add_title_text(slide, icon, 1, y, 0.6, 0.5, font_size=22, bold=False)
        add_title_text(slide, schema, 1.8, y, 3, 0.5, font_size=18, color=COLORS['primary'])
        add_body_text(slide, tables, 5, y, 7, 0.5, font_size=14)


def add_airbyte_slide(prs):
    """Add Airbyte slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔄 Data Ingestion: Airbyte", 0.5, 0.4, 12, 0.8, font_size=34)
    
    features = [
        ("🔌", "300+ Connectors", "Pre-built source/destination connectors"),
        ("🔄", "Incremental Sync", "Only sync changed data"),
        ("📊", "Schema Detection", "Auto-detect and evolve schemas"),
        ("🌐", "Open Source", "No vendor lock-in"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        y = 1.3 + i * 1.0
        add_title_text(slide, icon, 1, y, 0.6, 0.5, font_size=24, bold=False)
        add_title_text(slide, title, 1.8, y, 3, 0.5, font_size=20, color=COLORS['accent'])
        add_body_text(slide, desc, 1.8, y + 0.4, 10, 0.4, font_size=14)
    
    # Sync modes
    add_title_text(slide, "Sync Modes", 7.5, 1.3, 4, 0.5, font_size=18, color=COLORS['primary'])
    modes = ["• Full Refresh - Complete replacement", "• Incremental - Append new records", "• CDC - Change Data Capture"]
    for i, mode in enumerate(modes):
        add_body_text(slide, mode, 7.5, 1.9 + i * 0.5, 5, 0.4, font_size=13)


def add_dbt_overview_slide(prs):
    """Add dbt overview slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "⚙️ Data Transformation: dbt", 0.5, 0.4, 12, 0.8, font_size=34)
    add_body_text(slide, "data build tool - Transform data using SQL + software engineering best practices",
                  0.5, 1.1, 12, 0.4, font_size=18, color=COLORS['text'])
    
    features = [
        ("📝", "SQL-based", "Write transformations in SQL"),
        ("🔀", "DAG", "Automatic dependency management"),
        ("✅", "Testing", "Built-in data quality tests"),
        ("📚", "Documentation", "Auto-generated data catalog"),
        ("🔄", "Incremental", "Process only new data"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        y = 1.8 + i * 0.9
        add_title_text(slide, icon, 1, y, 0.6, 0.5, font_size=22, bold=False)
        add_title_text(slide, title, 1.8, y, 2.5, 0.5, font_size=18, color=COLORS['primary'])
        add_body_text(slide, desc, 4.5, y, 8, 0.5, font_size=14)


def add_dbt_layers_slide(prs):
    """Add dbt layers slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📊 dbt Model Layers", 0.5, 0.4, 12, 0.8, font_size=34)
    
    layers = """
    ┌───────────────────────────────────────────────────────────────┐
    │                      dbt Project Layers                        │
    ├───────────────────────────────────────────────────────────────┤
    │                                                               │
    │  ┌───────────┐      ┌───────────────┐      ┌─────────────┐   │
    │  │  Staging  │      │ Intermediate  │      │    Marts    │   │
    │  │  (88 SQL) │ ──►  │   dim + fact  │ ──►  │ (5 tables)  │   │
    │  │           │      │               │      │             │   │
    │  │  stg_*    │      │  dim_customer │      │ mart_sales  │   │
    │  │  (views)  │      │  dim_product  │      │ mart_ops    │   │
    │  │           │      │  fact_sales   │      │ mart_hr     │   │
    │  └───────────┘      └───────────────┘      └─────────────┘   │
    │                                                               │
    │  Purpose:           Purpose:              Purpose:            │
    │  Clean raw data     Business entities     Analytics-ready     │
    │  Rename columns     Relationships         Denormalized        │
    │  Cast types         Calculated metrics    Self-service BI     │
    │                                                               │
    └───────────────────────────────────────────────────────────────┘"""
    
    add_code_box(slide, layers, 0.5, 1.1, 12.3, 5.5, font_size=10)


def add_dbt_marts_slide(prs):
    """Add dbt marts slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📈 dbt Analytics Marts", 0.5, 0.4, 12, 0.8, font_size=34)
    
    marts = [
        ("mart_sales", "Order line item", "Revenue, territories, trends"),
        ("mart_customer_analytics", "Customer", "CLV, segmentation, churn"),
        ("mart_product_analytics", "Product", "Profitability, inventory"),
        ("mart_operations", "Order (PO/WO)", "Vendor performance, production"),
        ("mart_employee_performance", "Employee/period", "Quotas, territories"),
    ]
    
    # Headers
    add_body_text(slide, "Mart", 1, 1.3, 4, 0.4, font_size=14, color=COLORS['primary'])
    add_body_text(slide, "Grain", 5.5, 1.3, 3, 0.4, font_size=14, color=COLORS['primary'])
    add_body_text(slide, "Key Use Cases", 8.5, 1.3, 4, 0.4, font_size=14, color=COLORS['primary'])
    
    for i, (mart, grain, uses) in enumerate(marts):
        y = 1.8 + i * 0.7
        add_body_text(slide, mart, 1, y, 4, 0.4, font_size=14, color=COLORS['cyan'])
        add_body_text(slide, grain, 5.5, y, 3, 0.4, font_size=13)
        add_body_text(slide, uses, 8.5, y, 4, 0.4, font_size=13)
    
    # Supported analytics
    add_title_text(slide, "Analytics Supported", 1, 5.3, 11, 0.5, font_size=16, color=COLORS['accent'])
    add_body_text(slide, "✅ CLV • ✅ RFM Segmentation • ✅ Churn Prediction • ✅ Inventory Optimization • ✅ Territory Performance",
                  1, 5.8, 11, 0.4, font_size=13, color=COLORS['accent'])


def add_streamlit_slide(prs):
    """Add Streamlit dashboard slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📊 Streamlit Dashboard", 0.5, 0.4, 12, 0.8, font_size=34)
    
    pages = [
        ("🏠", "Overview", "KPIs, revenue trends, geographic map"),
        ("🤖", "AI Assistant", "Natural language queries"),
        ("💰", "Sales & Revenue", "Territory performance, CLV"),
        ("📦", "Product & Inventory", "Stock levels, turnover"),
        ("👥", "Customer Analytics", "Segmentation, churn risk"),
        ("👔", "HR Analytics", "Employee performance, quotas"),
        ("⚙️", "Operations", "Vendor metrics, production"),
        ("🛡️", "OPAL Demo", "Authorization demonstration"),
    ]
    
    for i, (icon, page, desc) in enumerate(pages):
        col = i % 2
        row = i // 2
        x = 1 + col * 6
        y = 1.3 + row * 1.3
        
        add_title_text(slide, icon, x, y, 0.5, 0.5, font_size=24, bold=False)
        add_title_text(slide, page, x + 0.6, y, 2.5, 0.5, font_size=18, color=COLORS['text'])
        add_body_text(slide, desc, x + 0.6, y + 0.4, 5, 0.4, font_size=12)


def add_ai_assistant_slide(prs):
    """Add AI assistant slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🤖 AI Analytics Assistant", 0.5, 0.4, 12, 0.8, font_size=34)
    add_body_text(slide, "Natural Language to SQL powered by GPT-4 / Claude",
                  0.5, 1.1, 12, 0.4, font_size=18, color=COLORS['primary'])
    
    flow = """
    User: "What is our revenue by territory?"
              │
              ▼
    ┌──────────────────────────────────────┐
    │          AI Processing               │
    │  1. Parse natural language           │
    │  2. Match to schema context          │
    │  3. Generate SQL query               │
    │  4. Validate for safety              │
    │  5. Execute against database         │
    │  6. Auto-visualize results           │
    └──────────────────────────────────────┘
              │
              ▼
    SQL: SELECT territory_name, SUM(order_total)
         FROM mart_sales GROUP BY territory_name
              │
              ▼
    📊 [Auto-generated bar chart]"""
    
    add_code_box(slide, flow, 2, 1.8, 9, 4.8, font_size=11)


def add_ai_architecture_slide(prs):
    """Add AI architecture slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🧠 AI Assistant Architecture", 0.5, 0.4, 12, 0.8, font_size=34)
    
    components = [
        ("schema_context.py", "Load table schemas, metrics, column info"),
        ("sql_generator.py", "Build prompts, call LLM API, generate SQL"),
        ("sql_validator.py", "Validate safety: SELECT only, whitelist"),
        ("visualizer.py", "Auto-detect chart type, format results"),
    ]
    
    for i, (file, desc) in enumerate(components):
        y = 1.3 + i * 1.1
        add_body_text(slide, file, 1, y, 3, 0.4, font_size=16, color=COLORS['cyan'])
        add_body_text(slide, desc, 1, y + 0.4, 11, 0.4, font_size=14)
    
    # Auto-sync note
    add_title_text(slide, "🔄 Auto-Sync with dbt", 1, 5.5, 11, 0.5, font_size=18, color=COLORS['accent'])
    add_body_text(slide, "Running ./run_dbt.sh run auto-generates schema_ai.md and allowed_tables.json",
                  1, 6, 11, 0.4, font_size=14, color=COLORS['accent'])


def add_ai_security_slide(prs):
    """Add AI security slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔒 AI Safety & Security", 0.5, 0.4, 12, 0.8, font_size=34)
    
    rules = [
        ("SELECT Only", "No INSERT, UPDATE, DELETE, DROP"),
        ("Table Whitelist", "Only query allowed mart tables"),
        ("Auto LIMIT", "Prevent huge result sets (max 1000)"),
        ("No Injection", "Block SQL injection patterns"),
        ("Read-Only", "No DDL operations allowed"),
    ]
    
    for i, (rule, desc) in enumerate(rules):
        y = 1.3 + i * 0.9
        add_title_text(slide, "✓", 1, y, 0.4, 0.4, font_size=18, color=COLORS['accent'])
        add_title_text(slide, rule, 1.5, y, 3, 0.4, font_size=18, color=COLORS['text'])
        add_body_text(slide, desc, 4.5, y, 8, 0.4, font_size=14)


def add_why_policy_slide(prs):
    """Add why policy-based access control slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🔐 Why Policy-Based Access Control?", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Problems
    add_title_text(slide, "Traditional Problems", 1, 1.3, 5, 0.5, font_size=20, color=COLORS['danger'])
    problems = ["🔓 Hard-coded permissions", "😰 Scattered logic", "📝 No audit trail", "🔄 Changes require deploy"]
    for i, p in enumerate(problems):
        add_body_text(slide, p, 1.2, 1.9 + i * 0.6, 5, 0.4, font_size=14)
    
    # Benefits
    add_title_text(slide, "Policy-as-Code Benefits", 7, 1.3, 5, 0.5, font_size=20, color=COLORS['accent'])
    benefits = ["📝 Declarative - readable format", "🔐 Centralized - single source", "📊 Auditable - track decisions", "⚡ Real-time - no deploys"]
    for i, b in enumerate(benefits):
        add_body_text(slide, b, 7.2, 1.9 + i * 0.6, 5, 0.4, font_size=14, color=COLORS['accent'])


def add_opa_opal_slide(prs):
    """Add OPA/OPAL slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🛡️ OPA + OPAL Stack", 0.5, 0.4, 12, 0.8, font_size=34)
    
    arch = """
    ┌────────────────────────────────────────────────────┐
    │              Policy Management                      │
    ├────────────────────────────────────────────────────┤
    │                                                    │
    │  ┌──────────────┐                                  │
    │  │Git Repository│ ← Policies (Rego) + Data (JSON)  │
    │  └──────┬───────┘                                  │
    │         │                                          │
    │         ▼                                          │
    │  ┌──────────────┐                                  │
    │  │ OPAL Server  │ ← Watches changes, pushes updates│
    │  └──────┬───────┘                                  │
    │         │ WebSocket                                │
    │         ▼                                          │
    │  ┌──────────────┐                                  │
    │  │ OPAL Client  │                                  │
    │  │  + OPA Engine│ ← Evaluates policies             │
    │  └──────┬───────┘                                  │
    │         │                                          │
    │         ▼                                          │
    │  ┌──────────────┐                                  │
    │  │ Applications │ ← Streamlit, APIs...             │
    │  └──────────────┘                                  │
    └────────────────────────────────────────────────────┘"""
    
    add_code_box(slide, arch, 1.5, 1.1, 10, 5.6, font_size=10)


def add_roles_slide(prs):
    """Add roles slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "👥 User Roles", 0.5, 0.4, 12, 0.8, font_size=34)
    
    roles = [
        ("admin", "Full access"),
        ("data_engineer", "All data + write"),
        ("analyst", "Read marts, AI"),
        ("viewer", "Limited read"),
        ("executive", "All dashboards"),
        ("sales_manager", "Sales + customer"),
        ("hr_manager", "Employee only"),
        ("operations_manager", "Inventory + production"),
    ]
    
    for i, (role, desc) in enumerate(roles):
        col = i % 2
        row = i // 2
        x = 1 + col * 6
        y = 1.3 + row * 1.2
        
        add_title_text(slide, f"• {role}", x, y, 3, 0.5, font_size=16, color=COLORS['primary'])
        add_body_text(slide, desc, x + 0.3, y + 0.45, 5, 0.4, font_size=13)


def add_access_matrix_slide(prs):
    """Add access matrix slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📊 Dashboard Access Matrix", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Headers
    headers = ["Dashboard", "Admin", "Analyst", "Viewer", "Exec", "Sales", "HR", "Ops"]
    for i, h in enumerate(headers):
        x = 0.6 + i * 1.55
        add_body_text(slide, h, x, 1.2, 1.5, 0.4, font_size=11, 
                      color=COLORS['primary'] if i == 0 else COLORS['text'])
    
    # Data
    data = [
        ("💰 Sales", "✅", "✅", "✅", "✅", "✅", "❌", "❌"),
        ("👔 HR", "✅", "❌", "❌", "✅", "❌", "✅", "❌"),
        ("⚙️ Operations", "✅", "✅", "✅", "✅", "❌", "❌", "✅"),
        ("👥 Customer", "✅", "✅", "❌", "✅", "✅", "❌", "❌"),
        ("📦 Inventory", "✅", "✅", "❌", "✅", "❌", "❌", "✅"),
        ("🤖 AI", "✅", "✅", "❌", "✅", "❌", "❌", "❌"),
    ]
    
    for row_idx, row in enumerate(data):
        y = 1.7 + row_idx * 0.65
        for col_idx, cell in enumerate(row):
            x = 0.6 + col_idx * 1.55
            color = COLORS['text'] if col_idx == 0 else (
                COLORS['accent'] if cell == "✅" else COLORS['danger']
            )
            add_body_text(slide, cell, x, y, 1.5, 0.4, font_size=13, color=color)


def add_start_stack_slide(prs):
    """Add starting the stack slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🚀 Starting the Stack", 0.5, 0.4, 12, 0.8, font_size=34)
    
    code = """# One command to start everything
./start.sh

# What happens:
# Step 1: PostgreSQL (port 5432)
# Step 2: Streamlit (port 8501)
# Step 3: dbt-docs (port 8080)
# Step 4: SQL Server (port 1433)
# Step 5: AdventureWorks installation
# Step 6: Airbyte (port 8000)
# Step 7: OPAL/OPA (port 8181)

# Service URLs:
# 📊 PostgreSQL     localhost:5432
# 📈 Streamlit      http://localhost:8501
# 📚 dbt Docs       http://localhost:8080
# 🔄 Airbyte        http://localhost:8000
# 🔐 OPA            http://localhost:8181"""
    
    add_code_box(slide, code, 1, 1.2, 11, 5.5, font_size=12)


def add_commands_slide(prs):
    """Add useful commands slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "💻 Useful Commands", 0.5, 0.4, 12, 0.8, font_size=34)
    
    commands = [
        ("Docker", "docker-compose ps | docker-compose logs -f streamlit"),
        ("dbt", "./run_dbt.sh run | ./run_dbt.sh test | ./run_dbt.sh build"),
        ("Airbyte", "abctl local status | abctl local credentials"),
        ("OPAL", "./setup.sh status | ./setup.sh test"),
        ("Stop", "./stop.sh"),
    ]
    
    for i, (category, cmds) in enumerate(commands):
        y = 1.3 + i * 1.0
        add_title_text(slide, category, 1, y, 2, 0.5, font_size=18, color=COLORS['primary'])
        add_code_box(slide, cmds, 3, y, 9.5, 0.6, font_size=11)


def add_takeaways_slide(prs):
    """Add key takeaways slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🎯 Key Takeaways", 0.5, 0.4, 12, 0.8, font_size=34)
    
    takeaways = [
        ("1", "Modern Data Stack", "Modular, best-of-breed components; ELT > ETL"),
        ("2", "Data Pipeline", "Airbyte → PostgreSQL → dbt (staging → marts)"),
        ("3", "Analytics", "Streamlit dashboards + AI Assistant with auto-sync"),
        ("4", "Security", "OPA/OPAL policy-as-code with RBAC"),
    ]
    
    for i, (num, title, desc) in enumerate(takeaways):
        y = 1.2 + i * 1.35
        add_title_text(slide, num, 1, y, 0.5, 0.5, font_size=28, color=COLORS['primary'])
        add_title_text(slide, title, 1.7, y, 10, 0.5, font_size=22, color=COLORS['text'])
        add_body_text(slide, desc, 1.7, y + 0.5, 10, 0.5, font_size=15)


def add_best_practices_slide(prs):
    """Add best practices slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "✅ Best Practices", 0.5, 0.4, 12, 0.8, font_size=34)
    
    practices = [
        ("dbt", "staging → intermediate → marts pattern"),
        ("dbt", "Document all models with descriptions"),
        ("AI", "Always validate generated SQL"),
        ("Security", "Default deny, least privilege"),
        ("Security", "Keep policies in version control"),
        ("Ops", "Use Docker Compose for consistency"),
    ]
    
    for i, (area, practice) in enumerate(practices):
        col = i % 2
        row = i // 2
        x = 1 + col * 6
        y = 1.3 + row * 1.3
        
        add_body_text(slide, area, x, y, 1.5, 0.4, font_size=14, color=COLORS['primary'])
        add_body_text(slide, practice, x, y + 0.4, 5.5, 0.4, font_size=13, color=COLORS['text'])


def add_next_steps_slide(prs):
    """Add next steps slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "🚀 Next Steps", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Immediate
    add_title_text(slide, "Immediate", 1, 1.3, 3, 0.5, font_size=18, color=COLORS['accent'])
    immediate = ["□ Deploy to your environment", "□ Connect your data sources", "□ Customize dashboards"]
    for i, item in enumerate(immediate):
        add_body_text(slide, item, 1.2, 1.9 + i * 0.5, 5, 0.4, font_size=14)
    
    # Short-term
    add_title_text(slide, "Short-term", 1, 3.6, 3, 0.5, font_size=18, color=COLORS['primary'])
    short = ["□ Add dbt models for your data", "□ Create custom AI prompts", "□ Define access policies"]
    for i, item in enumerate(short):
        add_body_text(slide, item, 1.2, 4.2 + i * 0.5, 5, 0.4, font_size=14)
    
    # Long-term
    add_title_text(slide, "Long-term", 7, 1.3, 3, 0.5, font_size=18, color=COLORS['purple'])
    long = ["□ Scale to production (K8s)", "□ Add more data sources", "□ Data quality monitoring", "□ CI/CD for dbt models"]
    for i, item in enumerate(long):
        add_body_text(slide, item, 7.2, 1.9 + i * 0.5, 5, 0.4, font_size=14)


def add_resources_slide(prs):
    """Add resources slide."""
    slide = add_slide(prs)
    add_bg(slide, prs)
    
    add_title_text(slide, "📚 Resources", 0.5, 0.4, 12, 0.8, font_size=34)
    
    # Documentation
    add_title_text(slide, "Project Docs", 1, 1.3, 5, 0.5, font_size=18, color=COLORS['primary'])
    docs = ["/README.md - Overview", "/dbt/README.md - dbt guide", "/opal/README.md - Security"]
    for i, doc in enumerate(docs):
        add_body_text(slide, f"• {doc}", 1.2, 1.9 + i * 0.5, 5, 0.4, font_size=13)
    
    # External
    add_title_text(slide, "External Resources", 1, 3.6, 5, 0.5, font_size=18, color=COLORS['cyan'])
    external = ["docs.getdbt.com", "docs.airbyte.com", "openpolicyagent.org"]
    for i, e in enumerate(external):
        add_body_text(slide, f"• {e}", 1.2, 4.2 + i * 0.5, 5, 0.4, font_size=13)
    
    # Try it
    add_title_text(slide, "Try It", 7, 1.3, 5, 0.5, font_size=18, color=COLORS['accent'])
    try_code = """git clone <repo>
cd data_warehouse
./start.sh
# Open http://localhost:8501"""
    add_code_box(slide, try_code, 7, 1.9, 5, 1.8, font_size=12)


def add_thank_you_slide(prs):
    """Add thank you slide."""
    slide = add_slide(prs)
    add_bg(slide, prs, RgbColor(23, 37, 84))
    
    add_title_text(slide, "🙏", 6, 1.5, 1.5, 1, font_size=72, bold=False)
    add_title_text(slide, "Thank You!", 3.5, 2.8, 6, 1, font_size=52)
    
    add_title_text(slide, "📊 Data-Driven Decisions Made Easy", 2.5, 4.2, 8, 0.6, 
                   font_size=26, color=COLORS['primary'])
    
    add_body_text(slide, "Questions?", 5.5, 5.5, 3, 0.5, font_size=24, color=COLORS['text'])


if __name__ == "__main__":
    print("Generating Data Warehouse Full Stack presentation...")
    
    prs = create_presentation()
    output_file = "Data_Warehouse_Full_Stack.pptx"
    prs.save(output_file)
    
    print(f"✅ Presentation saved to: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
