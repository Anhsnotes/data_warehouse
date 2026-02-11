#!/usr/bin/env python3
"""
Generate PowerPoint presentation for OPAL Authorization Workshop.
Run: python generate_presentation.py
Output: OPAL_Authorization_Workshop.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Helper to create RGB color
def RgbColor(r, g, b):
    """Create RGB color tuple for use with pptx."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Color scheme - Dark theme
COLORS = {
    'bg_dark': RgbColor(15, 23, 42),       # Slate 900
    'bg_card': RgbColor(30, 41, 59),       # Slate 800
    'primary': RgbColor(99, 102, 241),     # Indigo 500
    'accent': RgbColor(16, 185, 129),      # Emerald 500
    'text': RgbColor(248, 250, 252),       # Slate 50
    'text_muted': RgbColor(148, 163, 184), # Slate 400
    'danger': RgbColor(239, 68, 68),       # Red 500
    'warning': RgbColor(245, 158, 11),     # Amber 500
}


def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_section_slide(prs, "Part 1", "Introduction & Overview", "5 minutes")
    add_what_we_build_slide(prs)
    add_components_slide(prs)
    add_why_auth_slide(prs)
    add_section_slide(prs, "Part 2", "Architecture Deep Dive", "10 minutes")
    add_architecture_slide(prs)
    add_auth_flow_slide(prs)
    add_section_slide(prs, "Part 3", "OPAL & OPA Fundamentals", "10 minutes")
    add_opa_slide(prs)
    add_opal_slide(prs)
    add_rego_slide(prs)
    add_section_slide(prs, "Part 4", "Policy Design & Implementation", "10 minutes")
    add_user_model_slide(prs)
    add_role_permissions_slide(prs)
    add_access_matrix_slide(prs)
    add_rbac_policy_slide(prs)
    add_section_slide(prs, "Part 5", "Live Demo & Code Walkthrough", "10 minutes")
    add_demo_slide(prs)
    add_code_client_slide(prs)
    add_section_slide(prs, "Part 6", "Hands-on Exercise", "10 minutes")
    add_exercise_slide(prs)
    add_section_slide(prs, "Part 7", "Q&A and Wrap-up", "5 minutes")
    add_takeaways_slide(prs)
    add_best_practices_slide(prs)
    add_resources_slide(prs)
    add_thank_you_slide(prs)
    
    return prs


def add_slide(prs, layout_index=6):
    """Add a blank slide."""
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)


def add_title_text(slide, text, left, top, width, height, font_size=44, bold=True, color=None):
    """Add a title text box."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text']
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_body_text(slide, text, left, top, width, height, font_size=18, color=None):
    """Add body text."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS['text_muted']
    return shape


def add_code_box(slide, code, left, top, width, height, font_size=12):
    """Add a code box with monospace font."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(30, 30, 46)
    shape.line.color.rgb = RgbColor(69, 71, 90)
    
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    
    for i, line in enumerate(code.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Consolas'
        p.font.color.rgb = RgbColor(166, 227, 161)
    
    return shape


def add_title_slide(prs):
    """Add the title slide."""
    slide = add_slide(prs)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    # Shield icon
    add_title_text(slide, "🛡️", 5.5, 1.5, 2, 1.2, font_size=80, bold=False)
    
    # Main title
    add_title_text(slide, "OPAL Authorization", 2, 2.8, 9, 1, font_size=54)
    add_title_text(slide, "in Data Warehouse", 2, 3.6, 9, 0.8, font_size=42, 
                   color=COLORS['primary'])
    
    # Subtitle
    add_body_text(slide, "A Comprehensive 1-Hour Workshop", 2, 4.8, 9, 0.5, 
                  font_size=24, color=COLORS['text_muted'])
    
    # Footer
    add_body_text(slide, "Role-Based Access Control • Policy-as-Code • Real-time Updates",
                  2, 6.2, 9, 0.4, font_size=16, color=COLORS['text_muted'])


def add_agenda_slide(prs):
    """Add the agenda slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "📋 Agenda", 0.5, 0.4, 5, 0.8, font_size=36)
    
    agenda_items = [
        ("0:00 - 0:05", "Introduction & Overview", "🎯"),
        ("0:05 - 0:15", "Architecture Deep Dive", "🏗️"),
        ("0:15 - 0:25", "OPAL & OPA Fundamentals", "📚"),
        ("0:25 - 0:35", "Policy Design & Implementation", "⚙️"),
        ("0:35 - 0:45", "Live Demo & Code Walkthrough", "💻"),
        ("0:45 - 0:55", "Hands-on Exercise", "🔧"),
        ("0:55 - 1:00", "Q&A and Wrap-up", "❓"),
    ]
    
    for i, (time, topic, icon) in enumerate(agenda_items):
        y = 1.4 + i * 0.75
        
        # Time
        add_body_text(slide, time, 1, y, 2, 0.5, font_size=16, color=COLORS['accent'])
        
        # Icon + Topic
        add_body_text(slide, f"{icon}  {topic}", 3.5, y, 8, 0.5, font_size=20, 
                      color=COLORS['text'])


def add_section_slide(prs, part, title, duration):
    """Add a section divider slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(30, 27, 75)  # Indigo 950
    bg.line.fill.background()
    
    add_body_text(slide, part, 0.5, 2.5, 12, 0.6, font_size=24, color=COLORS['primary'])
    add_title_text(slide, title, 0.5, 3.2, 12, 1, font_size=48)
    add_body_text(slide, f"⏱️ {duration}", 0.5, 4.5, 3, 0.5, font_size=18, 
                  color=COLORS['text_muted'])


def add_what_we_build_slide(prs):
    """Add 'What We're Building' slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🏗️ What We're Building", 0.5, 0.4, 10, 0.8, font_size=36)
    
    add_body_text(slide, "A modern data warehouse with enterprise-grade authorization",
                  0.5, 1.2, 12, 0.5, font_size=20, color=COLORS['text'])
    
    architecture = """
┌─────────────────────────────────────────────────────────────┐
│                    Streamlit Dashboard                       │
│  ┌─────────┐ ┌─────────┐ ┌─────────┐ ┌─────────┐           │
│  │  Sales  │ │   HR    │ │   Ops   │ │   AI    │           │
│  └────┬────┘ └────┬────┘ └────┬────┘ └────┬────┘           │
│       └───────────┴───────────┴───────────┘                  │
│                         │                                    │
│                    ┌────▼────┐                               │
│                    │  OPAL   │  ◄── Policy-based Access     │
│                    └────┬────┘                               │
└─────────────────────────┼───────────────────────────────────┘
                          │
                    ┌─────▼─────┐
                    │ PostgreSQL │
                    └───────────┘"""
    
    add_code_box(slide, architecture, 1.5, 2, 10, 4.5, font_size=11)


def add_components_slide(prs):
    """Add project components slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🧩 Project Components", 0.5, 0.4, 10, 0.8, font_size=36)
    
    components = [
        ("📊", "Data Pipeline", "dbt + PostgreSQL", "Transform raw data into analytics marts"),
        ("🖥️", "Visualization", "Streamlit", "Interactive dashboards"),
        ("🔐", "Authorization", "OPAL + OPA", "Fine-grained access control"),
        ("🤖", "AI Assistant", "OpenAI/Anthropic", "Natural language queries"),
        ("💾", "Data Source", "AdventureWorks", "Sample enterprise data"),
    ]
    
    for i, (icon, name, tech, desc) in enumerate(components):
        y = 1.3 + i * 1.1
        
        # Icon
        add_title_text(slide, icon, 0.8, y, 0.8, 0.6, font_size=32, bold=False)
        
        # Name
        add_title_text(slide, name, 1.8, y, 3, 0.5, font_size=22, color=COLORS['text'])
        
        # Tech
        add_body_text(slide, tech, 5, y, 2.5, 0.5, font_size=16, color=COLORS['primary'])
        
        # Description
        add_body_text(slide, desc, 7.8, y, 5, 0.5, font_size=14, color=COLORS['text_muted'])


def add_why_auth_slide(prs):
    """Add 'Why Authorization Matters' slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "❓ Why Authorization Matters", 0.5, 0.4, 10, 0.8, font_size=36)
    
    # Problem section
    add_title_text(slide, "The Problem", 0.8, 1.4, 5, 0.5, font_size=24, color=COLORS['danger'])
    
    problems = [
        "Different users need different data access levels",
        "Compliance requirements (GDPR, SOC2, HIPAA)",
        "Prevent data breaches and unauthorized access",
        "Audit trail requirements",
    ]
    
    for i, problem in enumerate(problems):
        add_body_text(slide, f"•  {problem}", 1, 2 + i * 0.5, 5.5, 0.4, font_size=16)
    
    # Solution section
    add_title_text(slide, "The Solution", 7, 1.4, 5, 0.5, font_size=24, color=COLORS['accent'])
    
    solutions = [
        "Centralized policy management",
        "Real-time policy updates",
        "Declarative, auditable policies",
        "Separation of policy from code",
    ]
    
    for i, solution in enumerate(solutions):
        add_body_text(slide, f"✓  {solution}", 7.2, 2 + i * 0.5, 5, 0.4, font_size=16, 
                      color=COLORS['accent'])


def add_architecture_slide(prs):
    """Add system architecture slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🏛️ System Architecture", 0.5, 0.4, 10, 0.8, font_size=36)
    
    arch_diagram = """
                              ┌─────────────────────┐
                              │   Policy Repository │
                              │   (Git/File System) │
                              └──────────┬──────────┘
                                         │
                              ┌──────────▼──────────┐
                              │    OPAL Server      │
                              │  (Policy Admin)     │
                              └──────────┬──────────┘
                                         │
      ┌──────────────────────────────────┼──────────────────────────────────┐
      │                                  │                                   │
      ▼                                  ▼                                   ▼
┌───────────┐                    ┌─────────────┐                    ┌───────────┐
│OPAL Client│                    │ OPAL Client │                    │OPAL Client│
│  + OPA #1 │                    │   + OPA #2  │                    │  + OPA #3 │
└─────┬─────┘                    └──────┬──────┘                    └─────┬─────┘
      │                                 │                                 │
      ▼                                 ▼                                 ▼
┌───────────┐                    ┌─────────────┐                    ┌───────────┐
│ Streamlit │                    │ API Service │                    │    ETL    │
└───────────┘                    └─────────────┘                    └───────────┘"""
    
    add_code_box(slide, arch_diagram, 0.3, 1.2, 12.7, 5.5, font_size=10)


def add_auth_flow_slide(prs):
    """Add authorization flow slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🔄 Authorization Flow", 0.5, 0.4, 10, 0.8, font_size=36)
    
    steps = [
        ("1️⃣", "User Request", "User attempts to access a dashboard or data"),
        ("2️⃣", "Auth Check", "Application sends authorization query to OPA"),
        ("3️⃣", "Policy Evaluation", "OPA evaluates request against Rego policies"),
        ("4️⃣", "Decision", "OPA returns allow/deny decision"),
        ("5️⃣", "Enforcement", "Application shows/hides content based on decision"),
    ]
    
    for i, (num, title, desc) in enumerate(steps):
        y = 1.3 + i * 1.1
        
        add_title_text(slide, num, 0.8, y, 0.8, 0.6, font_size=28, bold=False)
        add_title_text(slide, title, 1.8, y, 3, 0.5, font_size=20, color=COLORS['accent'])
        add_body_text(slide, desc, 1.8, y + 0.4, 10, 0.4, font_size=14)


def add_opa_slide(prs):
    """Add OPA explanation slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🔓 What is OPA?", 0.5, 0.4, 10, 0.8, font_size=36)
    add_body_text(slide, "Open Policy Agent - A general-purpose policy engine",
                  0.5, 1.1, 12, 0.5, font_size=20, color=COLORS['primary'])
    
    features = [
        ("🔓", "Decoupled", "Separates policy from application code"),
        ("📝", "Declarative", "Policies written in Rego language"),
        ("⚡", "Fast", "Sub-millisecond policy evaluation"),
        ("🔌", "Universal", "Works with any application"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        y = 2 + i * 0.9
        add_title_text(slide, icon, 1, y, 0.6, 0.5, font_size=24, bold=False)
        add_title_text(slide, title, 1.8, y, 2, 0.5, font_size=18, color=COLORS['text'])
        add_body_text(slide, desc, 4, y, 8, 0.5, font_size=14)
    
    add_title_text(slide, "Use Cases", 7.5, 1.8, 4, 0.5, font_size=20, color=COLORS['accent'])
    use_cases = ["Kubernetes admission control", "API authorization", 
                 "Data filtering", "Feature flags"]
    for i, uc in enumerate(use_cases):
        add_body_text(slide, f"•  {uc}", 7.5, 2.4 + i * 0.5, 5, 0.4, font_size=14)


def add_opal_slide(prs):
    """Add OPAL explanation slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "📡 What is OPAL?", 0.5, 0.4, 10, 0.8, font_size=36)
    add_body_text(slide, "Open Policy Administration Layer - Policy distribution & data sync",
                  0.5, 1.1, 12, 0.5, font_size=20, color=COLORS['primary'])
    
    features = [
        ("🔄", "Real-time Updates", "Push policy changes instantly to all OPA instances"),
        ("📊", "Data Synchronization", "Keep authorization data in sync across clusters"),
        ("🔗", "Git Integration", "Policies as code with version control"),
        ("📨", "Event-driven", "WebSocket pub/sub for efficient updates"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        y = 2 + i * 1.1
        add_title_text(slide, icon, 1, y, 0.6, 0.5, font_size=24, bold=False)
        add_title_text(slide, title, 1.8, y, 3, 0.5, font_size=18, color=COLORS['text'])
        add_body_text(slide, desc, 1.8, y + 0.4, 10, 0.4, font_size=14)


def add_rego_slide(prs):
    """Add Rego language slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "📜 The Rego Language", 0.5, 0.4, 10, 0.8, font_size=36)
    
    code = """# Package declaration
package datawarehouse.authz

# Default deny - secure by default
default allow := false

# Rule: Allow if user has matching permission
allow if {
    some role in user_roles
    some permission in data.roles[role].permissions
    permission_matches(permission, input.action, input.resource)
}

# Helper: Check permission match
permission_matches(permission, action, resource) if {
    permission.action == action
    permission.resource == resource
}"""
    
    add_code_box(slide, code, 0.5, 1.2, 7, 5.5, font_size=12)
    
    # Key concepts on the right
    add_title_text(slide, "Key Concepts", 8, 1.2, 4, 0.5, font_size=20, color=COLORS['accent'])
    
    concepts = [
        ("Rules", "Boolean expressions"),
        ("Input", "Request data"),
        ("Data", "External data"),
        ("Sets", "Unique collections"),
    ]
    
    for i, (concept, desc) in enumerate(concepts):
        y = 1.9 + i * 0.7
        add_body_text(slide, f"{concept}:", 8, y, 2, 0.4, font_size=14, color=COLORS['primary'])
        add_body_text(slide, desc, 9.5, y, 3, 0.4, font_size=14)


def add_user_model_slide(prs):
    """Add user model slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "👥 User & Role Model", 0.5, 0.4, 10, 0.8, font_size=36)
    
    code = """{
  "users": {
    "admin@company.com": {
      "name": "System Administrator",
      "roles": ["admin"]
    },
    "senior.analyst@company.com": {
      "name": "Senior Data Analyst",
      "roles": ["analyst"]
    },
    "sales.manager.west@company.com": {
      "name": "Sales Manager - West",
      "roles": ["sales_manager"],
      "territories": ["1", "2", "3"]
    }
  }
}"""
    
    add_code_box(slide, code, 0.5, 1.2, 6.5, 5.5, font_size=13)
    
    # Explanation on the right
    add_title_text(slide, "Key Points", 7.5, 1.2, 5, 0.5, font_size=20, color=COLORS['accent'])
    
    points = [
        "Users identified by email",
        "Each user has one or more roles",
        "Additional attributes (territories, departments)",
        "Active/inactive status for access control",
    ]
    
    for i, point in enumerate(points):
        add_body_text(slide, f"•  {point}", 7.5, 1.9 + i * 0.6, 5, 0.4, font_size=14)


def add_role_permissions_slide(prs):
    """Add role permissions slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🔑 Role Permissions", 0.5, 0.4, 10, 0.8, font_size=36)
    
    code = """{
  "roles": {
    "admin": {
      "permissions": [
        {"action": "*", "resource": "*"}
      ]
    },
    "analyst": {
      "permissions": [
        {"action": "read", "resource": "mart_*"},
        {"action": "view", "resource": "dashboard.sales"},
        {"action": "view", "resource": "dashboard.ai_assistant"},
        {"action": "export", "resource": "mart_*"}
      ]
    },
    "viewer": {
      "permissions": [
        {"action": "read", "resource": "mart_*"},
        {"action": "view", "resource": "dashboard.sales"}
      ]
    }
  }
}"""
    
    add_code_box(slide, code, 0.5, 1.2, 7, 5.8, font_size=12)
    
    # Legend on the right
    add_title_text(slide, "Permission Types", 8, 1.2, 4, 0.5, font_size=20, color=COLORS['accent'])
    
    perms = [
        ("read", "Access table data"),
        ("view", "Access dashboards"),
        ("export", "Download data"),
        ("write", "Modify data"),
        ("*", "Wildcard (all)"),
    ]
    
    for i, (perm, desc) in enumerate(perms):
        y = 1.9 + i * 0.6
        add_body_text(slide, perm, 8, y, 1.5, 0.4, font_size=14, color=COLORS['primary'])
        add_body_text(slide, desc, 9.5, y, 3, 0.4, font_size=14)


def add_access_matrix_slide(prs):
    """Add access matrix slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "📊 Dashboard Access Matrix", 0.5, 0.4, 10, 0.8, font_size=36)
    
    # Header
    headers = ["Dashboard", "Admin", "Analyst", "Viewer", "Executive", "Sales", "HR", "Ops"]
    for i, h in enumerate(headers):
        x = 0.8 + i * 1.5
        add_body_text(slide, h, x, 1.3, 1.4, 0.4, font_size=12, 
                      color=COLORS['primary'] if i == 0 else COLORS['text'])
    
    # Data rows
    data = [
        ("💰 Sales", "✅", "✅", "✅", "✅", "✅", "❌", "❌"),
        ("👔 HR", "✅", "❌", "❌", "✅", "❌", "✅", "❌"),
        ("⚙️ Operations", "✅", "✅", "✅", "✅", "❌", "❌", "✅"),
        ("👥 Customer", "✅", "✅", "❌", "✅", "✅", "❌", "❌"),
        ("📦 Inventory", "✅", "✅", "❌", "✅", "❌", "❌", "✅"),
        ("🤖 AI Assistant", "✅", "✅", "❌", "✅", "❌", "❌", "❌"),
    ]
    
    for row_idx, row in enumerate(data):
        y = 1.8 + row_idx * 0.7
        for col_idx, cell in enumerate(row):
            x = 0.8 + col_idx * 1.5
            color = COLORS['text'] if col_idx == 0 else (
                COLORS['accent'] if cell == "✅" else COLORS['danger']
            )
            add_body_text(slide, cell, x, y, 1.4, 0.4, font_size=14, color=color)
    
    add_body_text(slide, "Each role has different access levels based on job function",
                  0.8, 6.2, 11, 0.4, font_size=14, color=COLORS['text_muted'])


def add_rbac_policy_slide(prs):
    """Add RBAC policy implementation slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "⚙️ RBAC Policy Implementation", 0.5, 0.4, 10, 0.8, font_size=36)
    
    code = """# Dashboard-specific access rules

# Sales dashboard - multiple roles allowed
allow if {
    input.action == "view"
    input.resource == "dashboard.sales"
    {"sales_manager", "executive", "analyst", "viewer"} & user_roles != set()
}

# HR dashboard - restricted access
allow if {
    input.action == "view"
    input.resource == "dashboard.hr"
    {"hr_manager", "executive"} & user_roles != set()
}

# AI Assistant - power users only
allow if {
    input.action == "view"
    input.resource == "dashboard.ai_assistant"
    {"analyst", "data_engineer", "executive"} & user_roles != set()
}"""
    
    add_code_box(slide, code, 0.5, 1.2, 12, 5.5, font_size=13)


def add_demo_slide(prs):
    """Add demo slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🖥️ Live Demo: OPAL Authorization Page", 0.5, 0.4, 12, 0.8, font_size=36)
    
    add_title_text(slide, "What We'll See:", 0.8, 1.3, 5, 0.5, font_size=24, color=COLORS['accent'])
    
    demos = [
        ("1️⃣", "Profile Switching", "Different users see different access"),
        ("2️⃣", "Dashboard Cards", "Visual access/denied indicators"),
        ("3️⃣", "Export Permissions", "Data export restrictions"),
        ("4️⃣", "Permissions Matrix", "Full access overview"),
    ]
    
    for i, (num, title, desc) in enumerate(demos):
        y = 2 + i * 0.9
        add_title_text(slide, num, 1, y, 0.6, 0.5, font_size=24, bold=False)
        add_title_text(slide, title, 1.8, y, 3, 0.5, font_size=18)
        add_body_text(slide, desc, 4.8, y, 7, 0.5, font_size=14)
    
    add_title_text(slide, "Key Observations:", 7, 1.3, 5, 0.5, font_size=24, color=COLORS['warning'])
    
    observations = [
        "Junior Analyst: Limited to Sales, Ops only",
        "Sales Manager: No HR or Operations access",
        "HR Director: Only sees HR dashboard",
        "Admin: Full access to everything",
    ]
    
    for i, obs in enumerate(observations):
        add_body_text(slide, f"•  {obs}", 7, 2 + i * 0.6, 6, 0.4, font_size=14)


def add_code_client_slide(prs):
    """Add code walkthrough slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "💻 Code: OPAL Client", 0.5, 0.4, 10, 0.8, font_size=36)
    
    code = """class OPALClient:
    def __init__(self, opa_url=None):
        self.opa_url = opa_url or 'http://localhost:8181'
    
    def authorize(self, user: str, action: str, resource: str) -> bool:
        \"\"\"Check authorization via OPA.\"\"\"
        response = httpx.post(
            f"{self.opa_url}/v1/data/datawarehouse/authz/allow",
            json={
                "input": {
                    "user": user,
                    "action": action,
                    "resource": resource
                }
            }
        )
        return response.json().get("result", False)

# Usage
client = OPALClient()
if client.authorize("analyst@co.com", "view", "dashboard.sales"):
    show_dashboard()"""
    
    add_code_box(slide, code, 0.5, 1.2, 12, 5.8, font_size=12)


def add_exercise_slide(prs):
    """Add exercise slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🔧 Hands-on Exercise", 0.5, 0.4, 10, 0.8, font_size=36)
    
    add_title_text(slide, "Task: Add a New Role", 0.8, 1.3, 10, 0.5, font_size=24, 
                   color=COLORS['accent'])
    
    add_body_text(slide, 'Create a "finance_analyst" role that can:', 0.8, 1.9, 11, 0.4, 
                  font_size=18, color=COLORS['text'])
    
    tasks = [
        "✓  View Sales and Customer Analytics dashboards",
        "✓  Export sales data only",
        "✗  Cannot access HR or Operations",
    ]
    
    for i, task in enumerate(tasks):
        color = COLORS['accent'] if task.startswith("✓") else COLORS['danger']
        add_body_text(slide, task, 1, 2.5 + i * 0.5, 10, 0.4, font_size=16, color=color)
    
    add_title_text(slide, "Steps:", 0.8, 4.2, 10, 0.5, font_size=20, color=COLORS['primary'])
    
    steps = [
        "1. Add role to opal/data/roles.json",
        "2. Add user to opal/data/users.json",
        "3. Update streamlit/pages/opal_demo.py",
        "4. Test the new role in the demo page",
    ]
    
    for i, step in enumerate(steps):
        add_body_text(slide, step, 1, 4.7 + i * 0.5, 10, 0.4, font_size=16)


def add_takeaways_slide(prs):
    """Add key takeaways slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "🎯 Key Takeaways", 0.5, 0.4, 10, 0.8, font_size=36)
    
    takeaways = [
        ("1️⃣", "Separation of Concerns", "Policy logic separate from application code"),
        ("2️⃣", "Declarative Policies", "Rego is expressive, readable, and auditable"),
        ("3️⃣", "Real-time Updates", "OPAL pushes policy changes instantly"),
        ("4️⃣", "Defense in Depth", "UI + server-side + row-level security"),
    ]
    
    for i, (num, title, desc) in enumerate(takeaways):
        y = 1.3 + i * 1.3
        add_title_text(slide, num, 1, y, 0.8, 0.6, font_size=28, bold=False)
        add_title_text(slide, title, 2, y, 10, 0.5, font_size=22, color=COLORS['accent'])
        add_body_text(slide, desc, 2, y + 0.5, 10, 0.4, font_size=16)


def add_best_practices_slide(prs):
    """Add best practices slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "✅ Best Practices", 0.5, 0.4, 10, 0.8, font_size=36)
    
    practices = [
        ("Default Deny", "Start with default allow := false"),
        ("Least Privilege", "Only grant necessary permissions"),
        ("Audit Logging", "Track all authorization decisions"),
        ("Test Policies", "Use OPA's built-in testing framework"),
        ("Version Control", "Keep policies in Git"),
        ("Fail Closed", "Deny access on errors"),
    ]
    
    for i, (practice, desc) in enumerate(practices):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 1.3 + row * 1.6
        
        add_title_text(slide, f"•  {practice}", x, y, 5.5, 0.5, font_size=20, color=COLORS['accent'])
        add_body_text(slide, desc, x + 0.3, y + 0.5, 5.5, 0.4, font_size=14)


def add_resources_slide(prs):
    """Add resources slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['bg_dark']
    bg.line.fill.background()
    
    add_title_text(slide, "📚 Resources & Next Steps", 0.5, 0.4, 10, 0.8, font_size=36)
    
    add_title_text(slide, "Documentation", 0.8, 1.3, 5, 0.5, font_size=22, color=COLORS['accent'])
    docs = [
        "• OPA Documentation: openpolicyagent.org/docs/",
        "• OPAL GitHub: github.com/permitio/opal",
        "• Rego Playground: play.openpolicyagent.org/",
    ]
    for i, doc in enumerate(docs):
        add_body_text(slide, doc, 1, 1.9 + i * 0.5, 6, 0.4, font_size=14)
    
    add_title_text(slide, "Project Files", 0.8, 3.6, 5, 0.5, font_size=22, color=COLORS['accent'])
    files = [
        "• opal/policies/rbac.rego - Main policy",
        "• opal/data/*.json - Users, roles, permissions",
        "• streamlit/pages/opal_demo.py - Demo page",
    ]
    for i, f in enumerate(files):
        add_body_text(slide, f, 1, 4.2 + i * 0.5, 6, 0.4, font_size=14)
    
    add_title_text(slide, "Try It Yourself", 7, 1.3, 5, 0.5, font_size=22, color=COLORS['accent'])
    code = """cd data_warehouse
./start.sh
# Navigate to http://localhost:8501
# Select "OPAL Authorization Demo\""""
    add_code_box(slide, code, 7, 1.9, 5.5, 2, font_size=13)


def add_thank_you_slide(prs):
    """Add thank you slide."""
    slide = add_slide(prs)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(30, 27, 75)
    bg.line.fill.background()
    
    add_title_text(slide, "🙏", 6, 1.5, 1.5, 1, font_size=72, bold=False)
    add_title_text(slide, "Thank You!", 3.5, 2.8, 6, 1, font_size=54)
    
    add_title_text(slide, "🛡️ Secure Data, Happy Users", 3, 4.2, 7, 0.6, 
                   font_size=28, color=COLORS['primary'])
    
    add_body_text(slide, "Questions?", 5.5, 5.5, 3, 0.5, font_size=24, color=COLORS['text'])


if __name__ == "__main__":
    print("Generating OPAL Authorization Workshop presentation...")
    
    prs = create_presentation()
    output_file = "OPAL_Authorization_Workshop.pptx"
    prs.save(output_file)
    
    print(f"✅ Presentation saved to: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
